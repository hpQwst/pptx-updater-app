#!/usr/bin/env python3
"""
Script para analisar estrutura de arquivo PPTX
Uso: python analyze_pptx.py <arquivo.pptx>
Retorna JSON com informações sobre slides, tabelas e gráficos
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print(json.dumps({"error": "python-pptx not installed"}), file=sys.stderr)
    sys.exit(1)

def analyze_pptx(pptx_path):
    """Analisa arquivo PPTX e retorna estrutura em JSON"""
    try:
        prs = Presentation(pptx_path)
        
        analysis = {
            "fileName": Path(pptx_path).name,
            "slideCount": len(prs.slides),
            "slides": []
        }
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slideNumber": slide_idx,
                "elements": [],
                "tables": [],
                "charts": []
            }
            
            for shape in slide.shapes:
                # Tabelas
                if shape.has_table:
                    table_data = {
                        "elementName": shape.name,
                        "elementType": "table",
                        "rowCount": len(shape.table.rows),
                        "colCount": len(shape.table.columns),
                        "preview": []
                    }
                    
                    # Extrair preview das primeiras 3 linhas
                    for row_idx, row in enumerate(shape.table.rows[:3]):
                        row_preview = [cell.text.strip() for cell in row.cells]
                        table_data["preview"].append(row_preview)
                    
                    slide_data["tables"].append(table_data)
                    slide_data["elements"].append({
                        "elementName": shape.name,
                        "elementType": "table"
                    })
                
                # Gráficos
                elif shape.has_chart:
                    chart_data = {
                        "elementName": shape.name,
                        "elementType": "chart",
                        "chartType": str(shape.chart.chart_type)
                    }
                    
                    slide_data["charts"].append(chart_data)
                    slide_data["elements"].append({
                        "elementName": shape.name,
                        "elementType": "chart"
                    })
                
                # Textos
                elif hasattr(shape, "text") and shape.text.strip():
                    text_preview = shape.text.replace('\n', ' ')[:100]
                    slide_data["elements"].append({
                        "elementName": shape.name,
                        "elementType": "text",
                        "content": text_preview
                    })
            
            analysis["slides"].append(slide_data)
        
        return analysis
    
    except Exception as e:
        return {"error": str(e)}

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: analyze_pptx.py <file.pptx>"}), file=sys.stderr)
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    result = analyze_pptx(pptx_file)
    print(json.dumps(result, ensure_ascii=False, indent=2))
