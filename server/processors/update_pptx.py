#!/usr/bin/env python3
"""
Script para atualizar arquivo PPTX com dados do Excel
Uso: python update_pptx.py <arquivo.pptx> <arquivo.xlsx> <mappings.json> <output.pptx>
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
    import openpyxl
except ImportError as e:
    print(json.dumps({"error": f"Missing dependency: {e}"}), file=sys.stderr)
    sys.exit(1)

def load_excel_data(excel_path, sheet_name, cell_range):
    """Carrega dados do Excel para um intervalo específico"""
    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        ws = wb[sheet_name]
        
        # Parse range like "A1:D10"
        parts = cell_range.split(":")
        if len(parts) != 2:
            return None
        
        start_cell = parts[0]
        end_cell = parts[1]
        
        # Extract coordinates
        start_col = ord(start_cell[0]) - ord('A')
        start_row = int(start_cell[1:]) - 1
        end_col = ord(end_cell[0]) - ord('A')
        end_row = int(end_cell[1:]) - 1
        
        data = []
        for row_idx in range(start_row, end_row + 1):
            row_data = []
            for col_idx in range(start_col, end_col + 1):
                cell = ws.cell(row=row_idx + 1, column=col_idx + 1)
                row_data.append(cell.value)
            data.append(row_data)
        
        return data
    except Exception as e:
        print(f"Error loading Excel data: {e}", file=sys.stderr)
        return None

def update_pptx_table(prs, slide_number, table_index, new_data):
    """Atualiza dados de uma tabela no PPTX"""
    try:
        slide = prs.slides[slide_number - 1]
        table_count = 0
        
        for shape in slide.shapes:
            if shape.has_table:
                if table_count == table_index:
                    table = shape.table
                    
                    # Update cells
                    for row_idx, row_data in enumerate(new_data):
                        if row_idx >= len(table.rows):
                            break
                        for col_idx, value in enumerate(row_data):
                            if col_idx >= len(table.columns):
                                break
                            cell = table.cell(row_idx, col_idx)
                            cell.text = str(value) if value is not None else ""
                    
                    return True
                
                table_count += 1
        
        return False
    except Exception as e:
        print(f"Error updating table: {e}", file=sys.stderr)
        return False

def process_pptx(pptx_path, excel_path, mappings, output_path):
    """Processa PPTX com mapeamentos do Excel"""
    try:
        prs = Presentation(pptx_path)
        updated_count = 0
        errors = []
        
        for mapping in mappings:
            try:
                element_type = mapping.get("elementType")
                
                if element_type == "table":
                    slide_number = mapping.get("slideNumber")
                    table_index = mapping.get("tableIndex", 0)
                    excel_sheet = mapping.get("excelSheet")
                    excel_range = mapping.get("excelRange")
                    
                    # Load data from Excel
                    new_data = load_excel_data(excel_path, excel_sheet, excel_range)
                    
                    if new_data and update_pptx_table(prs, slide_number, table_index, new_data):
                        updated_count += 1
                    else:
                        errors.append(f"Failed to update table on slide {slide_number}")
            
            except Exception as e:
                errors.append(f"Error processing mapping: {str(e)}")
        
        # Save updated PPTX
        prs.save(output_path)
        
        return {
            "success": True,
            "updatedCount": updated_count,
            "errors": errors,
            "outputPath": output_path
        }
    
    except Exception as e:
        return {
            "success": False,
            "error": str(e)
        }

if __name__ == "__main__":
    if len(sys.argv) < 5:
        print(json.dumps({"error": "Usage: update_pptx.py <pptx> <xlsx> <mappings.json> <output.pptx>"}), file=sys.stderr)
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    excel_file = sys.argv[2]
    mappings_file = sys.argv[3]
    output_file = sys.argv[4]
    
    # Load mappings
    with open(mappings_file, 'r') as f:
        mappings = json.load(f)
    
    result = process_pptx(pptx_file, excel_file, mappings, output_file)
    print(json.dumps(result, ensure_ascii=False, indent=2))
